

import os
import fitz  # PyMuPDF
from pptx import Presentation
from docx import Document
import tkinter as tk
from tkinter import filedialog
from MalwareDetector import MalwareDetector

from lists import Lists

class Process:
    def __init__(self):
        self.file_path = None  # Initialize file_path attribute to None  --> to get user select file path
        self.searchingLocation = []  # Initialize an empty list to ----->store search locations
        self.searchResult=[]            #----->store search result
        self.filePathtoDelete = []
        

#__________________________get user select file path   and    change name ---> file_path       
        # Define the setFilePath method
    def setFilePath(self, path):
        self.file_path = path
        #print("aa.py Selected path:", self.file_path)
    


# __________________ file_paths --> "blackbook.txt", "blackbook2.txt"  data line one by one pass MalwareDetector.add_signature()
    @staticmethod
    def read_signatures(file_paths):
 # Read malware signatures from files and add them to the detector    
         malware_detector = MalwareDetector()
         for file_path in file_paths:
            try:
                with open(file_path, 'r') as f:
                    for line in f:
                        signature = line.strip()

# ___  line by line get  for "blackbook.txt", "blackbook2.txt"  and pass line by line --->  malware_detector.add_signature(signature)                      
                        malware_detector.add_signature(signature)
            except Exception as e:
                print(  f"Error reading signature file {file_path}: {e}")
         return malware_detector


#@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@
# ______.extents methord with **os**   +each file paths and   +malweare signatures   to  compaire  
    @staticmethod
    def scan_text(content, malware_detector):

        #  scan text content for malware
        return malware_detector.is_infected(content)
    
#@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@
        

# ______.extents methord with **os** each file paths and malweare signatures   to  compaire 
#_______ get .pdf text and return Process.******scan_text(text, malware_detector)******
    @staticmethod
    def scan_pdf(file_path, malware_detector):
        # Scan a PDF file for malware
        try:
            doc = fitz.open(file_path)
            text = ""
            for page in doc:
                text += page.get_text()
            return Process.scan_text(text, malware_detector)
        except Exception as e:
            print(f"Error scanning PDF {file_path}: {e}")
            return False

# ______.extents methord with **os** each file paths and malweare signatures   to  compaire 
#_______ get .pptx text and return Process.******scan_text(text, malware_detector)******
    @staticmethod
    def scan_pptx(file_path, malware_detector):
        # Scan a PPTX file for malware
        try:
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text
            return Process.scan_text(text, malware_detector)
        except Exception as e:
            print(f"Error scanning PPTX {file_path}: {e}")
            return False

# ______.extents methord with **os** each file paths and malweare signatures   to  compaire  
#_______ get .docx text and return Process.******scan_text(text, malware_detector)******               
    @staticmethod
    def scan_docx(file_path, malware_detector):
        # Scan a DOCX file for malware
        try:
            doc = Document(file_path)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text
            return Process.scan_text(text, malware_detector)
        except Exception as e:
            print(f"Error scanning DOCX {file_path}: {e}")
            return False

# _____________________________pass the malweare signatures and **os** each file paths 
#  try to find OS each directly .extents  and pass the values in .extents methord with **os** each file paths and malweare signatures
    @staticmethod
    def scan_file(malware_detector, file_path):
        # Scan a file for malware based on its extension
        _, file_extension = os.path.splitext(file_path)
        if file_extension == ".pdf":
            return Process.scan_pdf(file_path, malware_detector)
        elif file_extension == ".pptx":
            return Process.scan_pptx(file_path, malware_detector)
        elif file_extension == ".docx":
            return Process.scan_docx(file_path, malware_detector)
        else:
            try:
                with open(file_path, 'r') as f:
                    content = f.read()
                return Process.scan_text(content, malware_detector)
            except Exception as e:
                print(f"Error scanning file {file_path}: {e}")
                return False
            

            
#______________________________main pass need to scan location , tree add done all signatores 
    def scan_directory(self, directory, malware_detector):
        # Scan a directory for malware
        malware_files = []

        try:
            print("Scanning files in progress:")
            for root, _, files in os.walk(directory):
                for idx, file in enumerate(files, start=1):
                    file_path = os.path.join(root, file)
                    print(f"Scanning file {idx}: {file_path}")

                    # Append the file_path to searchingLocation list
                    self.searchingLocation.append(file_path)
                    


#  -************************************************************** os path one by one  send 
                    if Process.scan_file(malware_detector, file_path):
                        # ---*****  append malweare file path to ----> malware_files[]
                        malware_files.append(file_path)
        except Exception as e:
            print(f"Error scanning directory {directory}: {e}")
        return malware_files


# --------- frontend give the parths array []  to deleted 
    def delete_file(self, file_path):
        # Delete a file
        try:
            os.unlink(file_path)
            print(f"Deleted: {file_path}")
        except Exception as e:
            print(f"Error deleting {file_path}: {e}")





#____ change name ---> file_path  after call --> main

    def main(self):
        if self.file_path:
            signature_files = ["blackbook.txt", "blackbook2.txt"]

#_____ send the signature_files to the       ---->  read_signatures(signature_files)   
            malware_detector = self.read_signatures(signature_files)
            print("\nScanning directory:", self.file_path)
            
            malware_files = self.scan_directory(self.file_path, malware_detector)
            if malware_files:
                print("\nMalware detected:")
                print("{:<70}{:<40}{:<10}".format("Location", "Name", "Size"))
                print("-" * 120)
                for file_path in malware_files:
                    file_name = os.path.basename(file_path)
                    file_size = os.path.getsize(file_path)
                    print("{:<70}{:<50}{:<20}".format(file_path, file_name, f"{file_size} bytes"))
                    self.searchResult.append(f"{file_path}      {file_name}     {file_size} bytes")
                    self.filePathtoDelete.append(file_path)
                print("-" * 120)
            else:
                print("No malware detected. You are safe.")
        else:
            print("No file path selected.")
            
if __name__ == "__main__":
    process = Process()
    process.main()